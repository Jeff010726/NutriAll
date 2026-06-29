from pathlib import Path
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "proposal"
IMG_DIR = OUT_DIR / "screenshots"
PPTX = OUT_DIR / "NutriAll_Diabetes_Care_网站设计方案.pptx"
OUT_DIR.mkdir(exist_ok=True)

W, H = 16, 9

COLORS = {
    "ink": RGBColor(18, 37, 29),
    "forest": RGBColor(31, 70, 51),
    "paper": RGBColor(246, 242, 234),
    "milk": RGBColor(255, 253, 248),
    "sage": RGBColor(221, 232, 215),
    "lime": RGBColor(209, 251, 143),
    "clay": RGBColor(217, 140, 109),
    "muted": RGBColor(96, 112, 105),
    "black": RGBColor(10, 18, 14),
}


def add_bg(slide, color="paper"):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(W), Inches(H))
    rect.fill.solid()
    rect.fill.fore_color.rgb = COLORS[color]
    rect.line.fill.background()
    return rect


def add_text(slide, text, x, y, w, h, size=24, color="ink", bold=False, font="Arial", align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    if align:
        p.alignment = align
    run = p.runs[0]
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = COLORS[color]
    return box


def add_label(slide, text, x, y, w=2.1, color="muted"):
    return add_text(slide, text.upper(), x, y, w, 0.28, size=9, color=color, bold=True)


def add_pill(slide, text, x, y, w, fill="forest", color="milk"):
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.42))
    pill.fill.solid()
    pill.fill.fore_color.rgb = COLORS[fill]
    pill.line.fill.background()
    add_text(slide, text, x + 0.16, y + 0.095, w - 0.32, 0.22, size=9, color=color, bold=True, align=PP_ALIGN.CENTER)
    return pill


def crop_image(path, out_name, ratio=16 / 9, focus_y=0.0):
    img = Image.open(path).convert("RGB")
    w, h = img.size
    current = w / h
    if current > ratio:
        new_w = int(h * ratio)
        left = (w - new_w) // 2
        box = (left, 0, left + new_w, h)
    else:
        new_h = int(w / ratio)
        top = int((h - new_h) * focus_y)
        top = max(0, min(top, h - new_h))
        box = (0, top, w, top + new_h)
    cropped = img.crop(box)
    out = OUT_DIR / "crops" / out_name
    out.parent.mkdir(exist_ok=True)
    cropped.save(out, quality=92)
    return out


def add_image(slide, path, x, y, w, h, radius=False):
    pic = slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
    return pic


def add_card(slide, x, y, w, h, fill="milk"):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS[fill]
    shape.line.color.rgb = RGBColor(226, 220, 210)
    shape.line.width = Pt(0.7)
    return shape


def title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, "paper")
    hero = crop_image(IMG_DIR / "home.png", "cover_home.jpg", ratio=1.08, focus_y=0.02)
    add_image(slide, hero, 9.0, 0.45, 5.95, 7.9)
    add_shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, Inches(0), Inches(0), Inches(5.8), Inches(3.25))
    add_shape.fill.solid()
    add_shape.fill.fore_color.rgb = COLORS["lime"]
    add_shape.line.fill.background()
    add_label(slide, "Website design proposal", 0.7, 0.68, 3.2)
    add_text(slide, "NutriAll\nDiabetes Care\n网站设计方案", 0.7, 1.25, 7.6, 2.6, size=44, color="forest", bold=True)
    add_text(slide, "将糖尿病营养管理业务从 NutriAll Wellness 中独立出来，构建一个更专业、更具转化力、更适合长期内容运营的品牌网站。", 0.75, 4.25, 6.55, 1.0, size=17, color="muted")
    add_pill(slide, "面向客户展示 / 中文提案", 0.75, 6.3, 2.2, "forest")
    add_text(slide, "2026", 0.75, 7.35, 1.2, 0.3, size=11, color="muted", bold=True)


def agenda_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, "forest")
    add_label(slide, "Proposal structure", 0.72, 0.65, 3, "lime")
    add_text(slide, "我们建议从“糖尿病单一业务”出发，重塑一个更清晰的线上增长入口。", 0.72, 1.22, 11.7, 1.2, size=32, color="milk", bold=True)
    items = [
        ("01", "项目背景", "为什么需要从综合营养站点中拆出糖尿病业务"),
        ("02", "品牌定位", "从服务介绍升级为专业的 diabetes care experience"),
        ("03", "页面体系", "首页、Program、Recipes、Research、Booking、Insurance、About"),
        ("04", "视觉方向", "高级、明亮、医疗可信，同时保留生活方式温度"),
        ("05", "转化路径", "保险权益、免费咨询、预约流程和内容运营闭环"),
    ]
    for i, (num, title, desc) in enumerate(items):
        y = 3.0 + i * 0.82
        add_text(slide, num, 0.85, y, 0.55, 0.3, size=12, color="lime", bold=True)
        add_text(slide, title, 1.55, y - 0.03, 2.2, 0.35, size=18, color="milk", bold=True)
        add_text(slide, desc, 3.7, y, 8.5, 0.32, size=13, color="sage")


def strategy_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, "paper")
    add_label(slide, "Strategy", 0.7, 0.62)
    add_text(slide, "不是做一个“服务页”，而是做一个可持续获客的糖尿病营养 care platform。", 0.7, 1.08, 10.7, 1.25, size=31, color="forest", bold=True)
    cards = [
        ("业务独立", "糖尿病用户需求明确，独立站点能建立更强专业心智。"),
        ("降低决策成本", "用户先理解服务、保险、食谱和研究，再进入预约。"),
        ("内容资产沉淀", "Recipes 与 Research 让网站具备长期 SEO 和教育价值。"),
        ("转化闭环", "Insurance benefits + Free call + Booking flow 形成清晰路径。"),
    ]
    for i, (title, body) in enumerate(cards):
        x = 0.75 + (i % 2) * 7.35
        y = 3.2 + (i // 2) * 2.0
        add_card(slide, x, y, 6.55, 1.45, "milk")
        add_text(slide, title, x + 0.35, y + 0.28, 2.5, 0.35, size=20, color="forest", bold=True)
        add_text(slide, body, x + 0.35, y + 0.72, 5.7, 0.45, size=13, color="muted")


def sitemap_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, "paper")
    add_label(slide, "Information architecture", 0.7, 0.62)
    add_text(slide, "建议网站结构", 0.7, 1.08, 4.5, 0.5, size=34, color="forest", bold=True)
    cols = [
        ("首页", ["品牌主张", "核心服务模块", "保险权益入口", "客户反馈 / 团队信任"]),
        ("Program", ["糖尿病营养计划", "Meal planning", "Glucose patterns", "Carb confidence"]),
        ("Recipes", ["食谱总览", "食谱详情", "原 NutriAll 食谱卡片", "糖尿病适配说明"]),
        ("Research", ["研究文章库", "通俗解读", "A1C / 饮食 / 运动", "持续内容运营"]),
        ("Booking", ["免费咨询", "保险/自费路径", "语言偏好", "套餐与联系方式"]),
    ]
    for i, (title, bullets) in enumerate(cols):
        x = 0.55 + i * 3.05
        add_card(slide, x, 2.2, 2.75, 5.55, "milk")
        add_text(slide, title, x + 0.25, 2.55, 2.2, 0.4, size=22, color="forest", bold=True)
        for j, b in enumerate(bullets):
            add_text(slide, f"- {b}", x + 0.25, 3.3 + j * 0.72, 2.25, 0.5, size=12, color="muted")


def visual_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, "paper")
    add_label(slide, "Visual language", 0.7, 0.62)
    add_text(slide, "高级、明亮、可信，但不冰冷。", 0.7, 1.05, 7.2, 0.7, size=36, color="forest", bold=True)
    add_text(slide, "视觉系统参考 Seed 的 premium wellness 语言：胶囊导航、full-width hero、错落图片拼贴、柔和动效、深绿与高亮绿形成品牌记忆点。", 0.72, 1.95, 8.4, 0.75, size=15, color="muted")
    swatches = [("Forest", "forest"), ("Paper", "paper"), ("Sage", "sage"), ("Lime", "lime"), ("Clay", "clay")]
    for i, (name, color) in enumerate(swatches):
        x = 0.75 + i * 1.35
        rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(3.15), Inches(0.92), Inches(0.92))
        rect.fill.solid()
        rect.fill.fore_color.rgb = COLORS[color]
        rect.line.fill.background()
        add_text(slide, name, x, 4.2, 1.2, 0.25, size=9, color="muted", bold=True, align=PP_ALIGN.CENTER)
    imgs = [
        (IMG_DIR / "home.png", "home_crop.jpg", 9.2, 0.8, 5.6, 3.0, 0.0),
        (IMG_DIR / "recipes.png", "recipes_crop.jpg", 8.0, 4.25, 3.5, 2.5, 0.0),
        (IMG_DIR / "booking.png", "booking_crop.jpg", 11.75, 4.25, 3.2, 2.5, 0.0),
    ]
    for p, name, x, y, w, h, fy in imgs:
        add_image(slide, crop_image(p, name, ratio=w / h, focus_y=fy), x, y, w, h)


def screenshot_slide(prs, title, subtitle, img_name, focus_y=0.0, dark=False):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, "forest" if dark else "paper")
    color = "milk" if dark else "forest"
    body = "sage" if dark else "muted"
    add_label(slide, "Page preview", 0.7, 0.55, 2.4, "lime" if dark else "muted")
    add_text(slide, title, 0.7, 0.98, 6.4, 0.95, size=31, color=color, bold=True)
    add_text(slide, subtitle, 0.72, 2.02, 5.7, 1.0, size=14, color=body)
    shot = crop_image(IMG_DIR / img_name, f"slide_{img_name}.jpg", ratio=16 / 8.15, focus_y=focus_y)
    add_image(slide, shot, 6.1, 0.6, 9.25, 7.75)


def split_showcase_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, "paper")
    add_label(slide, "Experience preview", 0.7, 0.55)
    add_text(slide, "多页面展示：从教育到转化", 0.7, 0.98, 7.2, 0.7, size=34, color="forest", bold=True)
    add_text(slide, "客户看到的不只是一个首页，而是一套可以承载内容、服务说明、保险解释和预约转化的完整网站体验。", 0.72, 1.75, 6.5, 0.75, size=14, color="muted")
    items = [
        ("Program", "program.png", 0.2),
        ("Recipes", "recipes.png", 0.0),
        ("Research", "research.png", 0.0),
        ("Booking", "booking.png", 0.0),
    ]
    for i, (label, img, fy) in enumerate(items):
        x = 0.7 + (i % 2) * 7.5
        y = 3.0 + (i // 2) * 2.6
        add_image(slide, crop_image(IMG_DIR / img, f"grid_{img}.jpg", ratio=3.05, focus_y=fy), x, y, 6.85, 2.05)
        add_pill(slide, label, x + 0.2, y + 0.18, 1.25, "forest")


def booking_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, "paper")
    add_label(slide, "Booking conversion", 0.7, 0.55)
    add_text(slide, "预约流程从“表单”升级为“引导式咨询入口”。", 0.7, 0.98, 8.3, 0.75, size=32, color="forest", bold=True)
    bullets = [
        "选择糖尿病目标：A1C、meal plan、glucose readings、prediabetes 等",
        "选择访问方式：线上、线下、免费 discovery call",
        "选择支付路径：保险权益、自费、不确定",
        "收集语言偏好与联系方式，便于人工 follow-up",
        "展示套餐价格与保险免责声明，减少后续沟通成本",
    ]
    for i, b in enumerate(bullets):
        add_text(slide, f"{i+1}. {b}", 0.85, 2.05 + i * 0.55, 6.0, 0.38, size=13, color="muted")
    add_image(slide, crop_image(IMG_DIR / "booking.png", "booking_flow_show.jpg", ratio=1.35, focus_y=0.2), 7.25, 0.7, 7.6, 7.6)


def roadmap_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, "forest")
    add_label(slide, "Roadmap", 0.7, 0.6, 2.4, "lime")
    add_text(slide, "从设计确认到正式上线的建议阶段", 0.7, 1.08, 8.0, 0.7, size=34, color="milk", bold=True)
    phases = [
        ("01", "方案确认", "页面结构、视觉风格、内容方向、转化路径确认"),
        ("02", "UI 深化", "桌面端/移动端关键页面设计，交互状态与组件规范"),
        ("03", "前端开发", "响应式页面、动效、内容数据、预约流程前端实现"),
        ("04", "内容与 SEO", "Recipes / Research / Insurance 内容完善和搜索优化"),
        ("05", "上线交付", "表单/CRM/预约系统接入，测试、部署和培训"),
    ]
    for i, (num, title, body) in enumerate(phases):
        x = 0.75 + i * 3.0
        add_card(slide, x, 3.0, 2.55, 3.2, "milk")
        add_text(slide, num, x + 0.22, 3.28, 0.7, 0.3, size=13, color="clay", bold=True)
        add_text(slide, title, x + 0.22, 3.76, 2.05, 0.35, size=20, color="forest", bold=True)
        add_text(slide, body, x + 0.22, 4.45, 2.05, 0.9, size=11, color="muted")


def closing_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, "paper")
    add_text(slide, "NutriAll Diabetes Care", 0.72, 0.75, 6.5, 0.5, size=20, color="forest", bold=True)
    add_text(slide, "下一步建议", 0.72, 1.75, 5.5, 0.65, size=38, color="forest", bold=True)
    add_text(slide, "确认业务定位、页面范围、保险信息表达边界与预约系统接入方式后，即可进入 UI 深化与正式开发阶段。", 0.75, 2.7, 6.8, 0.9, size=17, color="muted")
    add_pill(slide, "Design proposal ready", 0.75, 4.35, 2.25, "forest")
    add_image(slide, crop_image(IMG_DIR / "home.png", "closing_home.jpg", ratio=1.7, focus_y=0.0), 8.0, 0.6, 6.9, 7.7)


def build():
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    title_slide(prs)
    agenda_slide(prs)
    strategy_slide(prs)
    sitemap_slide(prs)
    visual_slide(prs)
    screenshot_slide(prs, "首页：建立专业糖尿病营养 care 心智", "首页承担品牌第一印象：明确糖尿病定位、展示 care areas、建立保险与预约入口。", "home.png", 0.0)
    screenshot_slide(prs, "Program：从流程说明升级为服务能力展示", "用 meal planning、glucose patterns、carb confidence 等模块解释服务价值，而不是普通 step card。", "program.png", 0.0, dark=True)
    screenshot_slide(prs, "Recipes：把原有食谱变成可浏览内容资产", "食谱总页与详情页承接 NutriAll 原内容，并加入糖尿病适配说明、营养信息和原始卡片。", "recipes.png", 0.0)
    screenshot_slide(prs, "Research：用通俗语言解释公开研究", "Research 页面用于承接长期内容运营和 SEO，让专业知识更容易被普通用户理解。", "research.png", 0.0, dark=True)
    booking_slide(prs)
    screenshot_slide(prs, "Insurance：把 $0 benefit 变成清晰转化路径", "不承诺绝对免费，而是解释 eligibility、benefit verification 和影响最终费用的条件。", "insurance.png", 0.0)
    split_showcase_slide(prs)
    roadmap_slide(prs)
    closing_slide(prs)
    prs.save(PPTX)


if __name__ == "__main__":
    build()
