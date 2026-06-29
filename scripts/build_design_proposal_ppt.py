from pathlib import Path
from zipfile import ZipFile

from PIL import Image, ImageDraw, ImageFilter
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
SHOT_DIR = ROOT / "proposal" / "screenshots"
OUT_DIR = ROOT / "proposal"
CROP_DIR = OUT_DIR / "ppt_crops"
OUT = OUT_DIR / "NutriAll_Diabetes_Care_网站设计方案_优化版.pptx"

W, H = 13.333, 7.5

COLORS = {
    "paper": "F6F2EA",
    "cream": "FFFDF8",
    "forest": "123B2A",
    "deep": "09251A",
    "sage": "879184",
    "sage_dark": "5C695E",
    "lime": "D2FB8D",
    "clay": "D98C6D",
    "rose": "F2C8B8",
    "mist": "E4EADF",
    "white": "FFFFFF",
    "ink": "163228",
}


def rgb(hex_color: str) -> RGBColor:
    hex_color = COLORS.get(hex_color, hex_color)
    hex_color = hex_color.strip("#")
    return RGBColor(int(hex_color[:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))


def path(name: str) -> Path:
    return SHOT_DIR / name


def crop(src: Path, name: str, top: int, height: int, width: int = 1440) -> Path:
    CROP_DIR.mkdir(parents=True, exist_ok=True)
    out = CROP_DIR / name
    with Image.open(src) as im:
        im = im.convert("RGB")
        h = im.height
        top = max(0, min(top, max(0, h - 1)))
        bottom = max(top + 1, min(top + height, h))
        box = (0, top, min(width, im.width), bottom)
        im.crop(box).save(out, quality=92)
    return out


def cover_crop(src: Path, name: str, aspect: float, focus_y: float = 0.2) -> Path:
    CROP_DIR.mkdir(parents=True, exist_ok=True)
    out = CROP_DIR / name
    with Image.open(src) as im:
        im = im.convert("RGB")
        iw, ih = im.size
        current = iw / ih
        if current > aspect:
            nw = int(ih * aspect)
            left = (iw - nw) // 2
            box = (left, 0, left + nw, ih)
        else:
            nh = int(iw / aspect)
            top = int((ih - nh) * focus_y)
            top = max(0, min(top, ih - nh))
            box = (0, top, iw, top + nh)
        im.crop(box).save(out, quality=92)
    return out


def make_soft_shadow(name: str, w_px: int, h_px: int, radius: int = 28) -> Path:
    CROP_DIR.mkdir(parents=True, exist_ok=True)
    out = CROP_DIR / name
    im = Image.new("RGBA", (w_px, h_px), (0, 0, 0, 0))
    draw = ImageDraw.Draw(im)
    draw.rounded_rectangle((20, 20, w_px - 30, h_px - 30), radius=radius, fill=(16, 35, 26, 90))
    im = im.filter(ImageFilter.GaussianBlur(18))
    im.save(out)
    return out


def add_rect(slide, x, y, w, h, fill, radius=True, line=None, transparency=0):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(fill)
    shp.fill.transparency = transparency
    if line:
        shp.line.color.rgb = rgb(line)
        shp.line.width = Pt(0.8)
    else:
        shp.line.fill.background()
    return shp


def add_text(slide, text, x, y, w, h, size=18, color="forest", bold=False, align="left",
             font="Microsoft YaHei", line_spacing=1.05):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for idx, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = rgb(color)
    return tb


def add_bullets(slide, items, x, y, w, h, size=14, color="ink", bullet_color="lime"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.space_after = Pt(7)
        p.line_spacing = 1.12
        r1 = p.add_run()
        r1.text = "• "
        r1.font.name = "Microsoft YaHei"
        r1.font.size = Pt(size + 1)
        r1.font.bold = True
        r1.font.color.rgb = rgb(bullet_color)
        r2 = p.add_run()
        r2.text = item
        r2.font.name = "Microsoft YaHei"
        r2.font.size = Pt(size)
        r2.font.color.rgb = rgb(color)
    return box


def add_label(slide, text, x, y, w=1.5, fill="lime", color="forest"):
    add_rect(slide, x, y, w, 0.34, fill, radius=True)
    add_text(slide, text, x + 0.12, y + 0.075, w - 0.24, 0.18, 8.5, color, True, "center")


def add_image(slide, img, x, y, w, h):
    return slide.shapes.add_picture(str(img), Inches(x), Inches(y), Inches(w), Inches(h))


def add_frame(slide, img, x, y, w, h, title=None, crop_focus=0.12):
    shadow = make_soft_shadow(f"shadow_{abs(hash((str(img), x, y, w, h))) % 100000}.png", 900, 560)
    slide.shapes.add_picture(str(shadow), Inches(x - 0.08), Inches(y - 0.02), Inches(w + 0.24), Inches(h + 0.25))
    add_rect(slide, x, y, w, h, "white", radius=True)
    if title:
        add_text(slide, title, x + 0.28, y + 0.16, w - 0.55, 0.22, 8.8, "sage_dark", True)
        for i, c in enumerate(["D98C6D", "E9D7AE", "879184"]):
            add_rect(slide, x + 0.14 + i * 0.12, y + 0.18, 0.055, 0.055, c, radius=True)
        pic_y = y + 0.44
        pic_h = h - 0.56
    else:
        pic_y = y + 0.08
        pic_h = h - 0.16
    fitted = cover_crop(Path(img), f"fit_{abs(hash((str(img), w, h, crop_focus))) % 100000}.jpg", w / pic_h, crop_focus)
    add_image(slide, fitted, x + 0.08, pic_y, w - 0.16, pic_h)


def add_phone_frame(slide, img, x, y, w, h, title=None, crop_focus=0.0):
    shadow = make_soft_shadow(f"phone_shadow_{abs(hash((str(img), x, y, w, h))) % 100000}.png", 420, 760, 52)
    slide.shapes.add_picture(str(shadow), Inches(x - 0.12), Inches(y + 0.02), Inches(w + 0.28), Inches(h + 0.28))
    add_rect(slide, x, y, w, h, "deep", radius=True)
    add_rect(slide, x + 0.08, y + 0.08, w - 0.16, h - 0.16, "white", radius=True)
    if title:
        add_text(slide, title, x, y - 0.3, w, 0.22, 9, "sage_dark", True, "center")
    fitted = cover_crop(Path(img), f"phone_fit_{abs(hash((str(img), w, h, crop_focus))) % 100000}.jpg", (w - 0.28) / (h - 0.34), crop_focus)
    add_image(slide, fitted, x + 0.14, y + 0.17, w - 0.28, h - 0.34)


def add_kicker(slide, text, x, y, color="sage_dark"):
    add_text(slide, text.upper(), x, y, 4.0, 0.22, 8.5, color, True)


def add_footer(slide, n):
    add_text(slide, f"{n:02d}", 12.55, 7.08, 0.35, 0.18, 8, "sage_dark", True, "right")
    add_text(slide, "NutriAll Diabetes Care · Website Design Proposal", 0.45, 7.08, 4.8, 0.18, 7.5, "sage_dark")


def add_title(slide, kicker, title, subtitle=None, x=0.6, y=0.52, w=5.9):
    add_kicker(slide, kicker, x, y)
    add_text(slide, title, x, y + 0.28, w, 1.25, 31, "forest", False, line_spacing=0.92)
    if subtitle:
        add_text(slide, subtitle, x, y + 1.62, w, 0.58, 12.5, "sage_dark", False, line_spacing=1.18)


def slide_bg(prs, fill="paper"):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, W, H, fill, radius=False)
    return s


def build_crops():
    crops = {
        "home_top": crop(path("home.png"), "home_top.jpg", 0, 920),
        "home_mid": crop(path("home.png"), "home_mid.jpg", 1220, 1120),
        "home_support": crop(path("home.png"), "home_support.jpg", 2050, 1050),
        "program_top": crop(path("program.png"), "program_top.jpg", 0, 960),
        "program_mid": crop(path("program.png"), "program_mid.jpg", 900, 1000),
        "recipes_top": crop(path("recipes.png"), "recipes_top.jpg", 0, 980),
        "recipes_grid": crop(path("recipes.png"), "recipes_grid.jpg", 850, 1100),
        "recipe_detail": crop(path("recipe-detail.png"), "recipe_detail.jpg", 0, 980),
        "research_top": crop(path("research.png"), "research_top.jpg", 0, 960),
        "research_detail": crop(path("research-detail.png"), "research_detail.jpg", 0, 980),
        "booking_top": crop(path("booking.png"), "booking_top.jpg", 0, 950),
        "booking_flow": crop(path("booking.png"), "booking_flow.jpg", 850, 1000),
        "insurance_top": crop(path("insurance.png"), "insurance_top.jpg", 0, 950),
        "about_top": crop(path("about.png"), "about_top.jpg", 0, 950),
        "mobile_home": crop(path("mobile-home.png"), "mobile_home.jpg", 0, 1600, width=390),
        "mobile_recipes": crop(path("mobile-recipes.png"), "mobile_recipes.jpg", 0, 1600, width=390),
        "mobile_research": crop(path("mobile-research.png"), "mobile_research.jpg", 0, 1600, width=390),
        "mobile_booking": crop(path("mobile-booking.png"), "mobile_booking.jpg", 0, 1600, width=390),
    }
    return crops


def build():
    crops = build_crops()
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)

    # 1 Cover
    s = slide_bg(prs, "paper")
    add_rect(s, 0, 0, W, H, "paper", radius=False)
    add_rect(s, 8.2, 0, 5.2, 7.5, "mist", radius=False)
    add_label(s, "CLIENT PROPOSAL", 0.72, 0.62, 1.55)
    add_text(s, "NutriAll\nDiabetes Care", 0.72, 1.13, 5.3, 1.55, 38, "forest", False, line_spacing=0.88)
    add_text(s, "糖尿病营养管理独立站\n网站设计方案", 0.74, 2.85, 5.15, 0.82, 18, "ink", True, line_spacing=1.16)
    add_text(s, "以专业营养服务、保险福利核验、\n公开研究解读和实用食谱为核心，\n建立一个可信、清晰、可转化的\n糖尿病护理入口。", 0.76, 4.05, 4.45, 1.18, 12.8, "sage_dark", False, line_spacing=1.18)
    add_text(s, "2026 · Website Concept Direction", 0.78, 6.72, 3.3, 0.22, 8.8, "sage_dark", True)
    add_frame(s, crops["home_top"], 5.78, 0.72, 6.85, 4.16, "Homepage Preview", 0)
    add_phone_frame(s, crops["mobile_home"], 10.95, 3.95, 1.42, 2.56, "Mobile", 0)
    add_frame(s, crops["recipes_top"], 7.1, 4.68, 2.25, 1.62, "Recipes", 0)
    add_frame(s, crops["booking_top"], 9.18, 4.35, 2.18, 1.78, "Booking", 0)
    add_footer(s, 1)

    # 2 Context
    s = slide_bg(prs)
    add_title(s, "PROJECT CONTEXT", "为什么需要一个糖尿病\n独立业务站", "将糖尿病业务从 NutriAll/NutriWell 的综合健康服务中拆分出来，需要更聚焦的叙事、更明确的转化路径和更强的信任建立。")
    add_bullets(s, [
        "目标用户不是泛健康人群，而是正在处理 A1C、血糖波动、饮食选择、保险福利和长期管理压力的人。",
        "网站需要从第一屏就明确：这是糖尿病营养管理，不是普通饮食建议。",
        "内容要兼顾专业感与可理解性：研究、食谱、保险、预约都必须用用户听得懂的话讲清楚。",
        "视觉上保留 NutriAll 的温度，同时做出更现代、更高级、更像独立品牌的表达。"
    ], 0.75, 2.72, 5.9, 2.35, 13.5)
    add_frame(s, crops["home_top"], 7.0, 0.72, 5.6, 3.2, "当前首页方向", 0)
    add_rect(s, 7.1, 4.38, 5.36, 1.55, "forest", radius=True)
    add_text(s, "核心定位", 7.45, 4.72, 1.4, 0.24, 10, "lime", True)
    add_text(s, "专业、可信、可访问的糖尿病营养护理入口", 7.45, 5.06, 4.4, 0.48, 20, "white", True)
    add_footer(s, 2)

    # 3 Architecture
    s = slide_bg(prs)
    add_title(s, "SITE ARCHITECTURE", "网站结构与路由规划", "结构以“了解服务 → 建立信任 → 浏览内容 → 核验保险 → 预约咨询”为主线，避免用户在关键转化前迷路。")
    cards = [
        ("首页", "品牌第一印象、核心价值、服务范围、保险提示、团队信任、预约入口"),
        ("Program", "糖尿病营养管理项目：饮食、血糖、A1C、碳水、文化饮食与长期陪伴"),
        ("Recipes", "糖尿病友好食谱总页与详情页，承接 SEO 和日常复访"),
        ("Research", "公开研究白话解读，建立科学可信度并降低理解门槛"),
        ("Insurance", "解释 $0 / 保险福利核验逻辑，减少价格焦虑"),
        ("Booking", "分步骤预约流程，把意向转为可跟进线索"),
        ("About", "团队、理念、服务语言和文化适配能力")
    ]
    x0, y0 = 0.72, 2.35
    for i, (name, desc) in enumerate(cards):
        col = i % 3
        row = i // 3
        x = x0 + col * 4.1
        y = y0 + row * 1.18
        add_rect(s, x, y, 3.72, 0.86, "cream", True, "mist")
        add_text(s, name, x + 0.22, y + 0.16, 1.25, 0.22, 13, "forest", True)
        add_text(s, desc, x + 1.08, y + 0.13, 2.35, 0.42, 8.5, "sage_dark", False, line_spacing=1.15)
    add_rect(s, 0.72, 5.36, 11.9, 0.82, "lime", True)
    add_text(s, "推荐客户沟通口径：先确认糖尿病独立站的定位与内容边界，再进入预约系统、保险核验和后端 CRM 的技术细节。", 1.02, 5.62, 11.2, 0.24, 12.2, "forest", True)
    add_footer(s, 3)

    # 4 Visual system
    s = slide_bg(prs)
    add_title(s, "VISUAL DIRECTION", "高级、清爽、可信，\n但不冰冷", "参考 Seed 的现代感与胶囊式导航交互，但糖尿病站点需要更明亮、更生活化、更强调真实护理场景。")
    add_frame(s, crops["home_top"], 6.7, 0.65, 5.8, 2.95, "Hero / Header", 0)
    add_frame(s, crops["home_support"], 6.7, 3.65, 5.8, 2.35, "Editorial Mosaic", 0.1)
    add_bullets(s, [
        "色彩：深绿建立医疗与营养可信度，亮绿色用于转化提示，米白保留温度，少量暖色提升亲和力。",
        "图片：从“咨询场景”扩展到食物、血糖工具、研究材料、团队、保险文件等多种主题，避免单调。",
        "排版：大标题负责记忆点，小段落负责快速理解，CTA 保持清晰、短句、可行动。",
        "交互：Header hover、下拉菜单、内容卡片展开和页面滚动都要有流畅、克制的动画。"
    ], 0.75, 2.5, 5.1, 2.55, 13)
    add_footer(s, 4)

    # 5 Homepage
    s = slide_bg(prs)
    add_title(s, "HOMEPAGE", "首页：把信任和转化\n放在第一屏", "首页要快速回答：你们做什么、我适不适合、下一步怎么开始。")
    add_frame(s, crops["home_top"], 6.55, 0.62, 5.82, 3.08, "首页首屏", 0)
    add_frame(s, crops["home_mid"], 6.55, 3.88, 2.82, 1.96, "服务模块", 0.15)
    add_frame(s, crops["home_support"], 9.55, 3.88, 2.82, 1.96, "内容矩阵", 0.1)
    add_bullets(s, [
        "首屏直接说明“Diabetes nutrition care”，不让用户猜测业务范围。",
        "保险福利提示放在顶部，降低价格门槛；文案保持合规，强调 eligible / benefits verification。",
        "中段用大面积横向内容矩阵展示支持范围，比传统 step list 更有品牌感。",
        "底部串联团队、研究、食谱、预约 CTA，形成完整的转化闭环。"
    ], 0.72, 2.92, 5.2, 2.25, 12.4)
    add_footer(s, 5)

    # 6 Navigation
    s = slide_bg(prs)
    add_title(s, "NAVIGATION SYSTEM", "统一 Header 与\nSeed 式下拉菜单", "导航不只是链接集合，而是一个轻量的内容导览系统。所有页面保持同一套导航内容，避免页面之间体验不一致。")
    add_frame(s, crops["home_top"], 0.72, 2.15, 5.55, 2.92, "胶囊导航 + Hero", 0)
    add_frame(s, crops["research_top"], 6.72, 2.15, 5.55, 2.92, "Research Mega Menu", 0)
    add_bullets(s, [
        "顶部状态更轻，滚动后进入胶囊形态，导航位置保持一致。",
        "菜单 hover 高亮以滑块方式移动，不是突然变色；下拉区域带轻微位移和透明度动画。",
        "Program / Recipes / Research 均展示真实内容入口：项目模块、单独菜谱、研究文章，而不是空泛分类。",
        "这套 Header 已模块化，后续新增页面可以统一继承。"
    ], 0.95, 5.45, 11.2, 0.88, 11.3)
    add_footer(s, 6)

    # 7 Program
    s = slide_bg(prs)
    add_title(s, "PROGRAM PAGE", "Program：把服务讲成\n用户能执行的支持", "Program 页面不是罗列服务，而是围绕糖尿病用户每天面对的真实决策组织内容。")
    add_frame(s, crops["program_top"], 6.65, 0.72, 5.72, 3.05, "Program 首屏", 0)
    add_frame(s, crops["program_mid"], 6.65, 4.02, 5.72, 1.88, "项目模块", 0.2)
    add_bullets(s, [
        "Meal planning：把“能吃什么”转成可重复执行的餐盘和购物建议。",
        "Blood sugar patterns：帮助用户理解数字背后的饮食、活动与时间因素。",
        "Carb confidence：不制造恐惧，强调份量、搭配和时机。",
        "Dietitian support：突出注册营养师陪伴，而不是一次性的资料下载。",
        "Culturally realistic meals：保留文化饮食，不把管理糖尿病等同于单一西式餐单。"
    ], 0.74, 2.42, 5.3, 2.7, 12.7)
    add_footer(s, 7)

    # 8 Recipes
    s = slide_bg(prs)
    add_title(s, "RECIPES SYSTEM", "Recipes：既是内容资产，\n也是复访入口")
    add_frame(s, crops["recipes_top"], 0.72, 2.15, 5.55, 2.82, "菜谱总页", 0)
    add_frame(s, crops["recipe_detail"], 6.78, 2.15, 5.55, 2.82, "菜谱详情页", 0)
    add_bullets(s, [
        "总页负责发现：按实用场景展示早餐、主餐、小食、甜点等内容。",
        "详情页负责执行：食材、步骤、营养提示、糖尿病友好说明都需要清晰可读。",
        "图片要明亮、真实、有食欲，同时避免过度“医疗化”。",
        "Header 下拉直接展示单独菜谱，减少从导航到具体内容的点击成本。"
    ], 0.92, 5.48, 11.1, 0.82, 11.5)
    add_footer(s, 8)

    # 9 Research
    s = slide_bg(prs)
    add_title(s, "RESEARCH LIBRARY", "Research：把公开研究\n翻译成普通人能懂的话", "研究页的价值不是堆论文，而是建立“这个品牌懂科学，也懂用户”的信任感。")
    add_frame(s, crops["research_top"], 6.62, 0.72, 5.82, 3.0, "Research Library", 0)
    add_frame(s, crops["research_detail"], 6.62, 4.0, 5.82, 1.88, "Article Detail", 0)
    add_bullets(s, [
        "每篇研究用白话标题、关键结论、适用人群、日常做法来呈现。",
        "主题覆盖 A1C、餐后步行、高纤维饮食、血糖监测、心血管风险等。",
        "文章详情页避免论文式表达，优先回答“这和我的饮食/运动/监测有什么关系”。",
        "这部分可以作为长期 SEO 与品牌权威建设的核心栏目。"
    ], 0.72, 2.52, 5.2, 2.3, 13)
    add_footer(s, 9)

    # 10 Booking
    s = slide_bg(prs)
    add_title(s, "BOOKING FLOW", "Booking：从兴趣到线索\n要足够顺滑", "预约页不是一张表单，而是一个低压力引导流程：先确认目标，再收集信息，最后给出下一步。")
    add_frame(s, crops["booking_top"], 6.62, 0.72, 5.82, 2.85, "预约首屏", 0)
    add_frame(s, crops["booking_flow"], 6.62, 3.84, 5.82, 2.1, "分步骤表单", 0.15)
    add_bullets(s, [
        "Step 1：用户目标，例如降低 A1C、建立餐单、看懂血糖读数。",
        "Step 2：就诊与诊断背景，帮助团队判断服务路径。",
        "Step 3：保险或自费偏好，提前处理价格与资格问题。",
        "Step 4：语言与沟通方式，强化文化适配。",
        "Step 5：联系方式与提交后说明，方便团队后续跟进。"
    ], 0.72, 2.88, 5.0, 2.25, 11.8)
    add_footer(s, 10)

    # 11 Insurance
    s = slide_bg(prs)
    add_title(s, "INSURANCE MESSAGE", "保险福利：降低门槛，\n但必须表达严谨")
    add_frame(s, crops["insurance_top"], 6.68, 0.82, 5.72, 3.18, "Insurance Page", 0)
    add_rect(s, 6.75, 4.38, 5.55, 1.42, "forest", True)
    add_text(s, "建议中文解释", 7.08, 4.7, 1.7, 0.22, 10, "lime", True)
    add_text(s, "符合条件的客户，可能通过保险福利支付 $0；最终费用取决于保险计划、资格核验和服务范围。", 7.08, 5.02, 4.75, 0.42, 13, "white", True)
    add_text(s, "用保险福利核验降低价格焦虑，但不能把 $0 表达成绝对承诺。", 0.72, 2.35, 5.2, 0.32, 12.3, "sage_dark")
    add_bullets(s, [
        "不要把 $0 写成绝对承诺，避免用户误解为所有人免费。",
        "页面需要解释：为什么营养咨询可能被保险覆盖、核验需要哪些信息、多久能得到结果。",
        "Booking 流程中应保留 insurance benefits 入口，和独立 Insurance 页面互相承接。",
        "客户提案阶段建议把这部分作为“转化亮点 + 合规待确认项”。"
    ], 0.72, 2.92, 5.35, 2.1, 12.4)
    add_footer(s, 11)

    # 12 About
    s = slide_bg(prs)
    add_title(s, "ABOUT + TEAM", "About：团队信任\n和品牌温度放在一起", "Team 不需要单独拆出一个弱页面，合并到 About 能让用户更快理解服务理念、团队背景和护理方式。")
    add_frame(s, crops["about_top"], 6.58, 0.72, 5.85, 3.22, "About / Team", 0)
    add_bullets(s, [
        "展示团队合照，承接 NutriAll 原站资产，增强真实感。",
        "突出注册营养师、文化饮食理解、多语言支持和长期陪伴。",
        "About 页也应有预约 CTA，而不是只做品牌介绍。",
        "如果后续客户有专家背书或资质证书，可加入信任徽章区。"
    ], 0.72, 2.5, 5.15, 2.1, 13)
    add_rect(s, 6.72, 4.42, 5.56, 1.1, "cream", True, "mist")
    add_text(s, "客户感知", 7.05, 4.72, 1.4, 0.22, 10, "sage_dark", True)
    add_text(s, "“这是一个真实团队在帮我管理糖尿病饮食，而不是一个自动生成的健康内容站。”", 8.18, 4.68, 3.65, 0.36, 12.2, "forest", True)
    add_footer(s, 12)

    # 13 Content plan
    s = slide_bg(prs)
    add_title(s, "CONTENT STRATEGY", "内容规划：让网站\n持续产生信任和流量", "糖尿病站点不能只靠首页转化。长期价值来自可持续的内容栏目和清晰的服务入口。")
    columns = [
        ("服务内容", ["糖尿病营养管理", "A1C 与血糖模式", "文化适配餐单", "保险福利说明"]),
        ("食谱内容", ["高蛋白早餐", "高纤维主食", "低糖甜点", "工作日备餐"]),
        ("研究解读", ["餐后步行", "膳食纤维", "碳水搭配", "心血管风险"]),
        ("转化内容", ["免费咨询", "保险核验", "服务套餐", "常见问题"])
    ]
    for i, (head, items) in enumerate(columns):
        x = 0.72 + i * 3.12
        add_rect(s, x, 2.36, 2.72, 2.65, "cream", True, "mist")
        add_text(s, head, x + 0.24, 2.72, 2.2, 0.25, 15, "forest", True)
        add_bullets(s, items, x + 0.24, 3.18, 2.28, 1.35, 10.4)
    add_rect(s, 0.72, 5.52, 11.9, 0.84, "forest", True)
    add_text(s, "SEO 方向：以“diabetes dietitian / diabetes nutrition / A1C diet / diabetic recipes / insurance covered nutrition counseling”等主题建立内容集群。", 1.05, 5.82, 11.0, 0.22, 11.8, "white", True)
    add_footer(s, 13)

    # 14 Responsive
    s = slide_bg(prs)
    add_title(s, "RESPONSIVE EXPERIENCE", "桌面端负责品牌感，\n移动端负责即时行动")
    add_frame(s, crops["home_top"], 0.8, 2.06, 5.05, 2.85, "Desktop Homepage", 0)
    add_phone_frame(s, crops["mobile_home"], 5.42, 1.35, 1.62, 3.12, "Home", 0)
    add_phone_frame(s, crops["mobile_booking"], 7.1, 1.35, 1.62, 3.12, "Booking", 0)
    add_phone_frame(s, crops["mobile_recipes"], 8.78, 1.35, 1.62, 3.12, "Recipes", 0)
    add_phone_frame(s, crops["mobile_research"], 10.46, 1.35, 1.62, 3.12, "Research", 0)
    add_bullets(s, [
        "桌面端展示完整品牌叙事、页面结构和专业感；移动端强调快速预约、食谱查看和文章阅读。",
        "手机端需要保留保险提示、Book Now、食谱入口和 Research 入口，不能只做简单堆叠。",
        "后续可按真实数据优化：移动端点击热区、表单掉点、文章阅读深度、食谱收藏率。"
    ], 0.95, 5.35, 11.25, 0.78, 10.8)
    add_footer(s, 14)

    # 15 Internationalization
    s = slide_bg(prs)
    add_title(s, "MULTILINGUAL PLAN", "多语言：服务扩张前\n先把内容体系设计好")
    add_rect(s, 0.78, 2.2, 3.2, 2.75, "cream", True, "mist")
    add_text(s, "第一阶段", 1.08, 2.55, 1.4, 0.22, 10, "clay", True)
    add_text(s, "中文 / 英文", 1.08, 2.95, 2.2, 0.32, 21, "forest", True)
    add_text(s, "先覆盖核心转化页面：Home、Program、Insurance、Booking、About。", 1.08, 3.55, 2.3, 0.72, 11, "sage_dark", False, line_spacing=1.18)
    add_rect(s, 4.32, 2.2, 3.2, 2.75, "cream", True, "mist")
    add_text(s, "第二阶段", 4.62, 2.55, 1.4, 0.22, 10, "clay", True)
    add_text(s, "内容本地化", 4.62, 2.95, 2.2, 0.32, 21, "forest", True)
    add_text(s, "食谱、研究文章、FAQ 根据不同文化饮食习惯重写，而不是逐字翻译。", 4.62, 3.55, 2.35, 0.78, 11, "sage_dark", False, line_spacing=1.18)
    add_rect(s, 7.86, 2.2, 3.2, 2.75, "cream", True, "mist")
    add_text(s, "第三阶段", 8.16, 2.55, 1.4, 0.22, 10, "clay", True)
    add_text(s, "语言偏好沉淀", 8.16, 2.95, 2.35, 0.32, 21, "forest", True)
    add_text(s, "预约与注册时收集语言偏好，后续 EDM、食谱提醒和活动通知按语言分发。", 8.16, 3.55, 2.32, 0.78, 11, "sage_dark", False, line_spacing=1.18)
    add_rect(s, 1.1, 5.52, 10.5, 0.75, "forest", True)
    add_text(s, "建议技术方向：所有页面文案进入可维护的内容字典 / CMS，避免未来每加一种语言都重做页面。", 1.42, 5.78, 9.85, 0.22, 11.5, "white", True)
    add_footer(s, 15)

    # 16 Login and CRM
    s = slide_bg(prs)
    add_title(s, "LOGIN + EMAIL CAPTURE", "登录系统：把访问者\n沉淀为可运营用户")
    add_rect(s, 0.78, 2.05, 4.1, 3.1, "forest", True)
    add_text(s, "用户注册", 1.12, 2.42, 1.5, 0.24, 11, "lime", True)
    add_text(s, "Email sign-up\n+ preference profile", 1.12, 2.88, 3.25, 0.74, 24, "white", True, line_spacing=1.05)
    add_text(s, "收集邮箱、语言偏好、糖尿病目标、保险/自费偏好、感兴趣的食谱类型。", 1.12, 4.06, 3.05, 0.58, 11.5, "mist", False, line_spacing=1.16)
    flows = [
        ("01", "食谱收藏", "用户保存喜欢的菜谱，系统可推荐相似内容。"),
        ("02", "研究订阅", "有新研究解读时按主题推送。"),
        ("03", "预约提醒", "未完成 Booking 的用户可触发 follow-up。"),
        ("04", "客户分层", "按目标、语言、保险状态分组运营。"),
    ]
    for i, (num, head, desc) in enumerate(flows):
        x = 5.35 + (i % 2) * 3.25
        y = 2.05 + (i // 2) * 1.55
        add_rect(s, x, y, 2.85, 1.16, "cream", True, "mist")
        add_text(s, num, x + 0.2, y + 0.22, 0.45, 0.2, 10, "clay", True)
        add_text(s, head, x + 0.72, y + 0.2, 1.65, 0.22, 13.5, "forest", True)
        add_text(s, desc, x + 0.72, y + 0.54, 1.85, 0.36, 8.8, "sage_dark", False, line_spacing=1.12)
    add_bullets(s, [
        "账号系统建议从轻量 Email 注册开始，不一开始做复杂会员体系。",
        "注册价值要明确：收藏食谱、订阅研究更新、接收优惠活动、查看预约进度。",
        "后续可连接 CRM / 邮件平台，把网站从展示页升级为客户运营入口。"
    ], 5.45, 5.2, 6.2, 0.8, 10.6)
    add_footer(s, 16)

    # 17 EDM automation
    s = slide_bg(prs)
    add_title(s, "EDM AUTOMATION", "EDM：用内容更新\n增加客户黏性")
    stages = [
        ("注册后", "欢迎邮件", "介绍服务、保险核验、推荐第一批食谱。"),
        ("每周", "食谱更新", "按早餐、主餐、低糖甜点等兴趣标签推送。"),
        ("每月", "研究摘要", "把新的糖尿病研究用普通语言总结给用户。"),
        ("活动期", "优惠提醒", "咨询活动、保险核验窗口、节日饮食专题。"),
        ("未预约", "转化唤醒", "提醒完成 Booking 或预约免费 discovery call。"),
    ]
    for i, (timing, head, desc) in enumerate(stages):
        x = 0.82 + i * 2.42
        add_rect(s, x, 2.35, 2.05, 2.55, "cream", True, "mist")
        add_text(s, timing, x + 0.18, 2.7, 0.8, 0.2, 9.5, "clay", True)
        add_text(s, head, x + 0.18, 3.1, 1.6, 0.24, 14.5, "forest", True)
        add_text(s, desc, x + 0.18, 3.58, 1.6, 0.62, 9.2, "sage_dark", False, line_spacing=1.15)
    add_rect(s, 0.82, 5.45, 11.45, 0.72, "lime", True)
    add_text(s, "关键指标：Email 注册率、预约转化率、食谱点击率、文章阅读率、活动邮件打开率、复访率。", 1.18, 5.7, 10.8, 0.2, 11.8, "forest", True)
    add_footer(s, 17)

    # 18 Roadmap
    s = slide_bg(prs)
    add_title(s, "IMPLEMENTATION ROADMAP", "落地路径建议", "为了控制客户决策成本，可以分阶段推进：先确认视觉和内容结构，再接入真实业务流程与增长系统。")
    steps = [
        ("01", "品牌与结构", "确认定位、页面结构、核心文案和视觉方向。"),
        ("02", "核心页面", "完成首页、Program、Recipes、Research、Insurance、Booking、About。"),
        ("03", "移动端 QA", "重点检查导航、CTA、表单、食谱与文章阅读体验。"),
        ("04", "系统接入", "预约表单、CRM、保险核验、邮件通知、分析工具。"),
        ("05", "增长升级", "多语言、Email 注册、EDM 自动化、内容订阅与用户分层。"),
    ]
    for i, (num, head, desc) in enumerate(steps):
        x = 0.78 + i * 2.48
        add_rect(s, x, 2.55, 2.05, 2.52, "cream", True, "mist")
        add_text(s, num, x + 0.18, 2.84, 0.58, 0.28, 16, "clay", True)
        add_text(s, head, x + 0.18, 3.32, 1.55, 0.24, 14, "forest", True)
        add_text(s, desc, x + 0.18, 3.82, 1.65, 0.74, 9.8, "sage_dark", False, line_spacing=1.18)
    add_footer(s, 18)

    # 19 Deliverables
    s = slide_bg(prs)
    add_title(s, "DELIVERABLES", "建议交付范围", "客户一旦确认方向，可以把交付拆成清晰包，便于报价与排期。")
    add_bullets(s, [
        "完整网站视觉系统：颜色、字体、按钮、Header、Footer、卡片、表单、CTA。",
        "页面设计与响应式实现：首页、Program、Recipes、Recipe Detail、Research、Research Detail、Insurance、Booking、About。",
        "内容结构：服务文案、研究文章模板、菜谱模板、保险说明、FAQ、预约表单字段。",
        "素材系统：NutriAll 原有资产整理、缺失视觉补齐、图片风格统一与压缩优化。",
        "转化与分析：预约提交、保险核验入口、事件追踪、基础 SEO。",
        "后续增长模块：多语言内容管理、Email 注册、用户偏好收集、EDM 自动推送。"
    ], 0.78, 2.3, 5.65, 2.8, 13.2)
    add_frame(s, crops["home_top"], 7.0, 1.0, 5.25, 2.46, "Design System", 0)
    add_frame(s, crops["recipes_grid"], 7.0, 3.78, 5.25, 2.05, "Content Templates", 0.15)
    add_footer(s, 19)

    # 20 Decisions
    s = slide_bg(prs)
    add_title(s, "CLIENT DECISIONS", "客户需要确认的事项", "这页用于会议收口：把美术方向讨论转成可执行的业务决策。")
    add_bullets(s, [
        "品牌命名：是否使用 NutriAll Diabetes Care，还是保留更独立的 diabetes brand。",
        "保险口径：$0 benefits 的准确表达、免责声明、核验流程和适用保险范围。",
        "预约流程：提交后由谁跟进、是否接入现有 CRM、是否需要自动邮件或短信。",
        "用户系统：是否先做 Email 注册、是否提供食谱收藏、是否需要用户偏好标签。",
        "EDM 运营：食谱更新、研究文章、优惠活动和预约提醒分别由谁维护。",
        "内容权限：NutriAll 原站食谱、团队合照、Logo 与图片资产是否全部可用于独立站。",
        "上线节奏：先上线核心转化页，还是同步完成研究库与食谱库。"
    ], 0.8, 2.18, 6.05, 3.3, 11.8)
    add_rect(s, 7.25, 2.05, 4.88, 2.95, "forest", True)
    add_text(s, "建议会议目标", 7.65, 2.58, 2.1, 0.22, 11, "lime", True)
    add_text(s, "让客户先对“糖尿病独立站的品牌方向与页面结构”点头，再进入报价、开发周期和后端集成讨论。", 7.65, 3.08, 3.8, 0.9, 19, "white", True, line_spacing=1.05)
    add_footer(s, 20)

    # 21 Closing
    s = slide_bg(prs, "forest")
    add_rect(s, 0, 0, W, H, "forest", radius=False)
    add_text(s, "NutriAll Diabetes Care", 0.82, 0.72, 5.2, 0.32, 13, "lime", True)
    add_text(s, "让糖尿病营养管理\n更容易开始，\n也更容易坚持。", 0.78, 1.58, 6.1, 1.9, 34, "white", False, line_spacing=0.95)
    add_text(s, "下一步：确认设计方向与内容范围后，\n可进入正式页面设计、内容整理、\n预约流程与保险核验系统规划。", 0.82, 5.38, 4.85, 0.78, 12.5, "mist", False, line_spacing=1.18)
    add_frame(s, crops["booking_top"], 6.95, 0.88, 5.2, 2.72, "Booking Direction", 0)
    add_frame(s, crops["research_top"], 7.5, 3.72, 4.1, 1.82, "Research / Recipes / Trust", 0)
    add_phone_frame(s, crops["mobile_booking"], 11.25, 3.58, 0.98, 1.82, None, 0)
    add_text(s, "21", 12.55, 7.08, 0.35, 0.18, 8, "lime", True, "right")

    prs.save(OUT)
    with ZipFile(OUT) as zf:
        media_count = len([n for n in zf.namelist() if n.startswith("ppt/media/")])
    return OUT, len(prs.slides), media_count


if __name__ == "__main__":
    out, slides, media = build()
    print(f"created={out}")
    print(f"slides={slides}")
    print(f"media={media}")
